from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()

# Use 16:9 aspect ratio
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Helper function to add title slide
def add_title_slide(prs, title_text, subtitle_text):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    # Set dark background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 23, 42) # Dark blue/slate
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(56, 189, 248)
    title.text_frame.paragraphs[0].font.size = Pt(64)
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text = subtitle_text
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(248, 250, 252)
    
# Helper function to add content slide
def add_content_slide(prs, title_text, content_list, bg_color=(15, 23, 42)):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*bg_color)
    
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(56, 189, 248)
    title.text_frame.paragraphs[0].font.bold = True
    
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = content_list[0]
    tf.paragraphs[0].font.color.rgb = RGBColor(248, 250, 252)
    tf.paragraphs[0].font.size = Pt(28)
    
    for item in content_list[1:]:
        p = tf.add_paragraph()
        p.text = item
        p.font.color.rgb = RGBColor(248, 250, 252)
        p.font.size = Pt(28)
        p.space_before = Pt(14)
        
    return slide

add_title_slide(prs, "OccluSynth", "Occlusion-Aware 3D Scene Reconstruction\nAditya Agarwal\nProblem Statement 09")

add_content_slide(prs, "The Problem: The Blind Spot", [
    "• Indoor robots only see what is in their line of sight.",
    "• Most geometry needed for safe motion is hidden behind furniture (Occlusions).",
    "• Classical fusion forces a binary choice:",
    "   - Assume unseen space is FREE -> Silent Collisions",
    "   - Assume unseen space is BLOCKED -> Paralysed Robot",
    "• Our Approach: Reconstruct an explicit third state: OCCLUDED."
])

add_content_slide(prs, "Data Flow & Pipeline", [
    "A 4-stage pipeline moving from raw pixels to a collision-safe trajectory:",
    "",
    "1. Perception: RGB -> Dense Depth (VGGT-Omega + RANSAC)",
    "2. Fusion: Visibility TSDF (4-State Voxel Grid)",
    "3. Completion: Predicting the Hidden (3D U-Net Completer)",
    "4. Planning: Risk-Graded A* (Uncertainty-Aware Map)"
])

add_content_slide(prs, "Technical Architecture", [
    "• Foundation: Meta's VGGT-Omega (feed-forward monocular depth).",
    "• Metric Grounding: RANSAC lifts relative depth to metric scale (~2.4% Error).",
    "• Visibility-Aware TSDF: We project rays into a 5cm voxel grid.",
    "• Voxel States:",
    "   - FREE (empty air)",
    "   - SURFACE (measured solid)",
    "   - OCCLUDED (unobserved, inferred)",
    "   - UNOBSERVABLE (no evidence, left alone)"
])

add_content_slide(prs, "Implementation: 3D U-Net & Planning", [
    "• The Completer: 14.7M parameter Encoder-Decoder architecture.",
    "• Masked L1 Loss: Penalizes errors only on known surfaces and occluded zones.",
    "• Uncertainty via MC Dropout: Model runs multiple times. Variance = Uncertainty.",
    "• Risk-Graded A* Planner: Applies a 'Risk Tax'.",
    "   - FREE cost = 1.0",
    "   - OCCLUDED cost = 1 + λp (Dynamic risk)",
    "   - SURFACE cost = Infinity"
])

add_content_slide(prs, "Open Source & AI Leveraged", [
    "• Foundation Models: Meta's VGGT-Omega (CVPR '25 Best Paper) used frozen.",
    "• Open Datasets:",
    "   - ScanNet v2: 1,513 indoor scenes (drove visibility fusion).",
    "   - 7-Scenes: Cross-dataset robustness check.",
    "• Agentic AI Workflow:",
    "   - Developed using LLM-assisted workflow (Claude Code).",
    "   - Automated literature tracking, ablation scripting, and 55+ test suite."
])

add_content_slide(prs, "Published Models & Results", [
    "• Published to Hugging Face:",
    "   - Repository: onemore-adi/occlusynth-completer",
    "   - 14.7M Param 3D U-Net (MIT License).",
    "• Real Results & Validation:",
    "   - Occluded MAE reduced from 45.3cm -> 27.1cm.",
    "   - Hidden Hazard Awareness improved from 0% -> 21%.",
    "   - Achieved closed-loop safe navigation bypassing occluded structural hazards."
])

prs.save('/Users/onemore_adi/OccluSynth/OccluSynth_Redesigned.pptx')
