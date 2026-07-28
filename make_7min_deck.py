#!/usr/bin/env python3
"""The 7-minute-format deck: title + 4 slides (problem / solution / innovation / KPI).

Flow on stage: slides 1-4 (~5 min with the talk) -> 2-min demo video
(demo_video/build/OccluSynth_hb2_2min.mp4) -> slide 5 (KPI + impact) to close.

Same Samsung landing-page design system as make_samsung_deck.py (which still
generates the full 12-slide deck — keep that one as the Q&A appendix).

Run:  <venv-with-python-pptx+Pillow>/python make_7min_deck.py
Out:  OccluSynth_7min.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "OccluSynth_7min_light.pptx")

# ---------------------------------------------------------------- design tokens
W, H = 13.333, 7.5
M = 0.8
CW = W - 2 * M

BLACK = "AvertaStd-Black"
BOLD = "AvertaStd-Bold"
BOOK = "AvertaStd-Regular"

BLUE = "1428A0"
AMBER = "E0A100"
AMBER_INK = "8A6200"
FREE = "2E7D32"
SURFACE = "C0272D"
UNOBS = "6E7681"

LIGHT = dict(bg="FFFFFF", ink="000000", ink2="55595F", ink3="8B9096",
             card="F5F6F7", card2="EDEEF0", amber_ink=AMBER_INK, hair="E4E5E7")
DARK = dict(bg="0B0E14", ink="FFFFFF", ink2="A8AEB5", ink3="767C84",
            card="161B23", card2="1E242E", amber_ink=AMBER, hair="252B34")

Y_KICKER = 0.62
Y_HEAD = 1.10
Y_DECK = 1.94
Y_BODY = 2.52
Y_BODY_NODECK = 2.10
Y_FOOT = 6.83


def rgb(h):
    return RGBColor.from_string(h)


def _no_outline(shape):
    shape.line.fill.background()
    shape.shadow.inherit = False


def rect(slide, x, y, w, h, fill, radius=0.06):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sh.adjustments[0] = min(0.5, radius / min(w, h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = rgb(fill)
    _no_outline(sh)
    return sh


def ring(slide, x, y, d, color, weight=1.25):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    sh.fill.background()
    sh.line.color.rgb = rgb(color)
    sh.line.width = Pt(weight)
    sh.shadow.inherit = False
    return sh


def dot(slide, x, y, d, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    sh.fill.solid()
    sh.fill.fore_color.rgb = rgb(color)
    _no_outline(sh)
    return sh


def tri(slide, x, y, size, color, rotation=90):
    sh = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                Inches(x), Inches(y), Inches(size * 0.85), Inches(size))
    sh.rotation = rotation
    sh.fill.solid()
    sh.fill.fore_color.rgb = rgb(color)
    _no_outline(sh)
    return sh


def _supsplit(text, ov):
    import re
    parts = re.split(r"\{([^}]*)\}", text)
    out = []
    for i, part in enumerate(parts):
        if not part:
            continue
        o = dict(ov)
        if i % 2 == 1:
            o["sup"] = True
        out.append((part, o))
    return out or [("", ov)]


def tbox(slide, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, spec in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = spec.get("align", PP_ALIGN.LEFT)
        if spec.get("space_after") is not None:
            p.space_after = Pt(spec["space_after"])
        if spec.get("line_spacing"):
            p.line_spacing = spec["line_spacing"]
        runs = spec["runs"]
        if isinstance(runs, str):
            runs = [(runs, {})]
        expanded = []
        for text, ov in runs:
            expanded.extend(_supsplit(text, ov))
        for text, ov in expanded:
            r = p.add_run()
            r.text = text
            f = r.font
            f.name = ov.get("font", spec.get("font", BOOK))
            f.size = Pt(ov.get("size", spec.get("size", 13)))
            f.bold = False
            f.color.rgb = rgb(ov.get("color", spec.get("color", "000000")))
            rPr = r._r.get_or_add_rPr()
            spc = ov.get("spc", spec.get("spc"))
            if spc:
                rPr.set("spc", str(int(spc * 100)))
            if ov.get("sup"):
                rPr.set("baseline", "30000")
    return tb


# ---------------------------------------------------------------- measurement
_FONT_FILES = {
    BLACK: os.path.expanduser("~/Library/Fonts/AvertaStd-Black.ttf"),
    BOLD: os.path.expanduser("~/Library/Fonts/AvertaStd-Bold.ttf"),
    BOOK: os.path.expanduser("~/Library/Fonts/AvertaStd-Regular.ttf"),
}
_cache = {}
WARNINGS = []


def measure(text, font, size):
    from PIL import ImageFont
    key = (font, round(size * 4))
    if key not in _cache:
        _cache[key] = ImageFont.truetype(_FONT_FILES[font], int(size * 4))
    return _cache[key].getlength(text) / 4.0


def wrap_lines(text, font, size, width_in):
    limit = width_in * 72.0
    lines, cur = 1, ""
    for word in text.split():
        trial = (cur + " " + word).strip()
        if measure(trial, font, size) <= limit or not cur:
            cur = trial
        else:
            lines += 1
            cur = word
    return lines


def fits(tag, text, font, size, width_in, height_in, leading=1.32):
    text = text.replace("{", "").replace("}", "")
    n = wrap_lines(text, font, size, width_in)
    need = n * size * leading / 72.0
    if need > height_in + 0.012:
        WARNINGS.append(f"{tag}: needs {need:.2f}in ({n} lines) in {height_in:.2f}in")
    return need


def pill(slide, x, y, text, fill, fg, size=10.5, font=BOLD, pad=0.30, h=0.34, spc=0.9):
    w = measure(text, font, size) / 72.0 + 2 * pad + (spc * len(text) / 72.0)
    sh = rect(slide, x, y, w, h, fill, radius=h / 2)
    tf = sh.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.color.rgb = rgb(fg)
    r._r.get_or_add_rPr().set("spc", str(int(spc * 100)))
    return sh, w


prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]


def new_slide(dark=False):
    s = prs.slides.add_slide(BLANK)
    t = DARK if dark else LIGHT
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = rgb(t["bg"])
    return s, t


def header(slide, t, kicker, headline, deck=None, head_size=38):
    pill(slide, M, Y_KICKER, kicker.upper(), BLUE, "FFFFFF")
    fits(f"head:{headline[:24]}", headline, BLACK, head_size, CW, 0.78)
    tbox(slide, M, Y_HEAD, CW, 0.8,
         [dict(runs=headline, font=BLACK, size=head_size, color=t["ink"])])
    if deck:
        fits(f"deck:{headline[:20]}", deck, BOOK, 14.5, CW - 1.2, 0.62)
        tbox(slide, M, Y_DECK, CW - 1.2, 0.66,
             [dict(runs=deck, font=BOOK, size=14.5, color=t["ink2"], line_spacing=1.3)])


def label(slide, t, x, y, w, text, color=None, size=10):
    tbox(slide, x, y, w, 0.22,
         [dict(runs=text.upper(), font=BOLD, size=size,
               color=color or t["ink3"], spc=1.2)])


# ================================================================= SLIDE 1 — title
s, t = new_slide(dark=True)
pill(s, M, 0.95, "PROBLEM STATEMENT 09  ·  FINAL PRESENTATION ROUND", BLUE, "FFFFFF")
tbox(s, M, 1.72, CW, 1.66,
     [dict(runs="OCCLUSYNTH", font=BLACK, size=92, color="FFFFFF", spc=2.2)])
tbox(s, M, 3.34, 9.6, 1.0,
     [dict(runs=[("Occlusion-aware 3D scene reconstruction in partially observable ", {}),
                 ("real-world environments", {"color": AMBER})],
           font=BOOK, size=20, color=DARK["ink2"], line_spacing=1.35)])

meta = [("Team", "onemore_adi"), ("Member", "Aditya Agarwal"),
        ("Institute", "National Institute of Technology, Rourkela"),
        ("Repository", "github.com/onemore-adi/OccluSynth")]
cols = [M, M + 2.05, M + 4.10, M + 7.70]
widths = [1.9, 1.9, 3.4, 3.9]
for (k, v), cx, cwid in zip(meta, cols, widths):
    label(s, t, cx, 5.42, cwid, k)
    tbox(s, cx, 5.72, cwid, 0.52,
         [dict(runs=v, font=BOOK, size=13.5, color="FFFFFF", line_spacing=1.25)])

# ================================================================= SLIDE 2 — problem
s, t = new_slide()
header(s, t, "The problem", "Robots can't see behind the couch.",
       "An indoor robot sees only line-of-sight. Classical fusion forces one dishonest "
       "choice about everything else — and both choices fail.")

# left column: the two big numbers
tbox(s, M, Y_BODY + 0.06, 3.3, 1.20,
     [dict(runs="61%", font=BLACK, size=64, color=t["ink"])])
tbox(s, M, Y_BODY + 1.06, 3.05, 0.75,
     [dict(runs="of a cluttered room's observable volume is occluded — hidden behind "
                "furniture.", font=BOOK, size=12, color=t["ink2"], line_spacing=1.3)])
tbox(s, M, Y_BODY + 1.92, 3.3, 1.20,
     [dict(runs="0%", font=BLACK, size=64, color=AMBER_INK)])
tbox(s, M, Y_BODY + 2.92, 3.05, 0.80,
     [dict(runs="of that hidden geometry is recovered by observation — structural: no "
                "sensor measures behind a surface.", font=BOOK, size=12, color=t["ink2"],
           line_spacing=1.3)])

RX, RW = M + 3.72, CW - 3.72
fail = [("Call the unseen space FREE", SURFACE, "Silent collisions",
         "The planner drives straight through the hidden chair leg it never measured. Unsafe."),
        ("Call the unseen space BLOCKED", UNOBS, "Paralysed robot",
         "Every occluded voxel becomes a wall. The robot freezes and never reaches the goal. Useless.")]
cy = Y_BODY
for chip, ccol, title, body in fail:
    rect(s, RX, cy, RW, 1.32, t["card"], radius=0.14)
    pill(s, RX + 0.34, cy + 0.26, chip.upper(), ccol, "FFFFFF", size=9, pad=0.24, h=0.29)
    tbox(s, RX + 0.34, cy + 0.66, 3.0, 0.34,
         [dict(runs=title, font=BOLD, size=17, color=t["ink"])])
    tbox(s, RX + 3.42, cy + 0.62, RW - 3.78, 0.66,
         [dict(runs=body, font=BOOK, size=12.5, color=t["ink2"], line_spacing=1.32)])
    cy += 1.50

rect(s, RX, cy + 0.06, RW, 1.06, "0B0E14", radius=0.14)
tbox(s, RX + 0.34, cy + 0.28, RW - 0.68, 0.7,
     [dict(runs=[("Warehouse AMRs, home robots, inspection bots — safety-critical indoor "
                  "autonomy cannot ship on ", {}),
                 ("“assume it's empty.”", {"font": BOLD, "color": "FFFFFF"})],
           font=BOOK, size=13.5, color=DARK["ink2"], line_spacing=1.35)])

# ================================================================= SLIDE 3 — solution
s, t = new_slide()
header(s, t, "Our solution", "An honest third state — reconstructed.",
       "OccluSynth labels every voxel with what the sensor actually knew, predicts only "
       "what is genuinely inferable, and prices the risk of every guess.")

# four-state legend row (compact)
states = [(FREE, "FREE", "seen empty"), (SURFACE, "SURFACE", "measured solid"),
          (AMBER, "OCCLUDED", "hidden — predicted"), (UNOBS, "UNOBSERVABLE", "left alone")]
sx = M
SLY = Y_BODY + 0.04
for col, name, desc in states:
    hero = name == "OCCLUDED"
    chip_w = 0.34
    rect(s, sx, SLY, chip_w, 0.34, col, radius=0.10)
    nm_w = measure(name, BOLD, 12.5) / 72.0
    tbox(s, sx + chip_w + 0.14, SLY - 0.035, nm_w + 0.12, 0.26,
         [dict(runs=name, font=BOLD, size=12.5,
               color=AMBER_INK if hero else t["ink"])])
    de_w = measure(desc, BOOK, 11) / 72.0
    tbox(s, sx + chip_w + 0.14, SLY + 0.20, de_w + 0.12, 0.24,
         [dict(runs=desc, font=BOOK, size=11, color=t["ink3"])])
    sx += chip_w + 0.14 + max(nm_w, de_w) + 0.55

# five-stage pipeline
stages = [("1", "Perception", "VGGT depth from RGB; RANSAC pins metric scale"),
          ("2", "Fusion", "rays carve a 4-state voxel grid"),
          ("3", "Completion", "3D U-Net predicts SDF in occluded space only"),
          ("4", "Uncertainty", "MC-dropout: confidence per voxel"),
          ("5", "Planner", "risk-graded A* prices every guess")]
sw, sgap = 2.05, 0.37
SY = Y_BODY + 0.78
for i, (num, name, desc) in enumerate(stages):
    x = M + i * (sw + sgap)
    hero = name == "Completion"
    rect(s, x, SY, sw, 1.86, t["card2"] if hero else t["card"], radius=0.16)
    dot(s, x + 0.28, SY + 0.26, 0.40, AMBER if hero else BLUE)
    tbox(s, x + 0.28, SY + 0.325, 0.40, 0.3,
         [dict(runs=num, font=BLACK, size=12.5,
               color="000000" if hero else "FFFFFF", align=PP_ALIGN.CENTER)])
    tbox(s, x + 0.28, SY + 0.78, sw - 0.56, 0.3,
         [dict(runs=name, font=BOLD, size=14, color=t["ink"])])
    tbox(s, x + 0.28, SY + 1.10, sw - 0.56, 0.72,
         [dict(runs=desc, font=BOOK, size=11, color=t["ink2"], line_spacing=1.25)])
    if i < len(stages) - 1:
        tri(s, x + sw + 0.135, SY + 0.86, 0.13, t["ink3"])

rect(s, M, SY + 2.10, CW, 1.00, "0B0E14", radius=0.14)
tbox(s, M + 0.42, SY + 2.30, CW - 0.84, 0.64,
     [dict(runs=[("The restraint is the trust: ", {"font": BOLD, "color": AMBER}),
                 ("we never guess outside the camera's view, and every amber voxel ships "
                  "with a confidence the planner can price.", {})],
           font=BOOK, size=13.5, color="FFFFFF", line_spacing=1.35)])

# ================================================================= SLIDE 4 — innovation
s, t = new_slide()
header(s, t, "Innovation", "Knows what it doesn't know.",
       "Against the closest published systems — Atlas, NeuralRecon, VGGT, DiffInDScene, "
       "Behind the Veil — OccluSynth is the first to do all five:")

caps = ["Reconstructs occluded geometry from a single pass",
        "Works from sparse metric anchors — no depth sensor",
        "Attaches a confidence to every predicted voxel",
        "Scores the occluded region as a safety problem — first open benchmark",
        "Closes the loop into a navigation planner"]
LW2 = 6.1
cy = Y_BODY + 0.10
for cap in caps:
    dot(s, M, cy + 0.055, 0.16, BLUE)
    tbox(s, M + 0.34, cy - 0.02, LW2 - 0.34, 0.55,
         [dict(runs=cap, font=BOLD, size=13.5, color=t["ink"], line_spacing=1.25)])
    cy += 0.585

RX3 = M + LW2 + 0.55
label(s, t, RX3, Y_BODY + 0.06, CW - LW2 - 0.55, "What makes it work")
work = [("Visibility mask as conditioning",
         "generation is fenced to the recoverable region — the model is never rewarded "
         "for guessing where it has no right to."),
        ("Recall-first by design",
         "a phantom obstacle costs a slowdown; a missed one costs a collision. The "
         "asymmetry is priced, per voxel."),
        ("Reproducible end to end",
         "open weights, 55+ automated tests, ScanNet-native benchmark — no simulator "
         "required.")]
wy = Y_BODY + 0.42
for wname, wdesc in work:
    rect(s, RX3, wy, CW - LW2 - 0.55, 0.94, t["card"], radius=0.13)
    tbox(s, RX3 + 0.28, wy + 0.14, CW - LW2 - 1.11, 0.70,
         [dict(runs=[(wname + " — ", {"font": BOLD, "color": t["ink"]}), (wdesc, {})],
               font=BOOK, size=11, color=t["ink2"], line_spacing=1.25)])
    wy += 1.06

rect(s, M, 6.42, CW, 0.62, t["card2"], radius=0.13)
tri(s, M + 0.34, 6.63, 0.16, BLUE)
tbox(s, M + 0.62, 6.55, CW - 1.0, 0.4,
     [dict(runs=[("Next — the demo: ", {"font": BOLD, "color": t["ink"]}),
                 ("the whole pipeline on a real apartment, RGB frames to completed mesh, "
                  "in two minutes.", {})],
           font=BOOK, size=12.5, color=t["ink2"])], anchor=MSO_ANCHOR.MIDDLE)

# ================================================================= SLIDE 5 — KPI (dark)
s, t = new_slide(dark=True)
header(s, t, "The numbers", "Measured, not promised.",
       "Aggregate over held-out ScanNet scenes the model never saw — every metric "
       "against its line-of-sight baseline.")

stats = [("Hidden surfaces recovered", "0%", "58%", "geometry no camera captured"),
         ("Hidden-hazard awareness", "0%", "21%", "anticipated before first sight"),
         ("Occluded F-score @ 5 cm", "0%", "32%", "TSDF cannot see behind walls"),
         ("Occluded MAE", "45.3", "27.1", "depth error in cm — nearly halved")]
scw, sgp = 2.77, 0.22
for i, (lab, base, big, note) in enumerate(stats):
    x = M + i * (scw + sgp)
    rect(s, x, Y_BODY, scw, 2.10, t["card"], radius=0.16)
    tbox(s, x + 0.30, Y_BODY + 0.28, scw - 0.60, 0.52,
         [dict(runs=lab.upper(), font=BOLD, size=9.5, color=t["ink3"], spc=1.0,
               line_spacing=1.25)])
    bw = measure(base, BLACK, 20) / 72.0
    tbox(s, x + 0.30, Y_BODY + 0.86, bw + 0.08, 0.4,
         [dict(runs=base, font=BLACK, size=20, color=t["ink3"])])
    tri(s, x + 0.30 + bw + 0.13, Y_BODY + 0.955, 0.12, t["ink3"])
    fits(f"kpi:{lab}", big, BLACK, 34, scw - 1.0 - bw, 0.62)
    tbox(s, x + 0.30 + bw + 0.38, Y_BODY + 0.74, scw - 0.72 - bw, 0.62,
         [dict(runs=big, font=BLACK, size=34, color=AMBER)])
    tbox(s, x + 0.30, Y_BODY + 1.56, scw - 0.60, 0.44,
         [dict(runs=note, font=BOOK, size=11, color=t["ink2"], line_spacing=1.25)])

tbox(s, M, 5.00, CW, 0.60,
     [dict(runs=[("Plus: ", {"font": BOLD, "color": "FFFFFF"}),
                 ("2.4% metric-depth error (10/10 scenes) · 15.5% collision-avoidance on a "
                  "13.6 m closed-loop path · 55+ automated tests · open weights, "
                  "reproducible benchmark.", {})],
           font=BOOK, size=12.5, color=t["ink2"], line_spacing=1.35)])

tbox(s, M, 5.72, CW, 0.42,
     [dict(runs="Warehouse AMRs. Home robots. Rescue drones. Every machine that must act "
                "before it can look.", font=BOLD, size=15, color="FFFFFF")])
tbox(s, M, 6.28, CW, 0.5,
     [dict(runs=[("The future of mapping isn't seeing more — ", {}),
                 ("it's knowing what you cannot see.", {"font": BOLD, "color": AMBER})],
           font=BOOK, size=16, color=t["ink2"])])

# ---------------------------------------------------------------- save
prs.save(OUT)
print(f"saved {OUT}  ({len(prs.slides._sldIdLst)} slides)")
if WARNINGS:
    print(f"{len(WARNINGS)} fit warnings:")
    for w in WARNINGS:
        print("  !", w)
else:
    print("no fit warnings")
