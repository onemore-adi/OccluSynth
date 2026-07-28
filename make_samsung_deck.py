#!/usr/bin/env python3
"""Rebuild OccluSynth_Final_Round.pptx in the Samsung landing-page design language.

Content is preserved from the original deck; only the design system changes.
Reference: samsung.com/us — SamsungSharpSans 700 headlines over generous whitespace,
SamsungOne body, near-monochrome ink, pill-shaped chips, brand blue #1428A0 used
sparingly, dark sections for premium/hero moments.

Local stand-ins (installed on this machine, same as the 5-min video):
  SamsungSharpSans -> AvertaStd-Black / AvertaStd-Bold
  SamsungOne       -> AvertaStd-Regular

Run:  scratch-venv/bin/python make_samsung_deck.py
Out:  OccluSynth_Final_Round_Samsung.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "OccluSynth_Final_Round_Samsung.pptx")

# ---------------------------------------------------------------- design tokens
W, H = 13.333, 7.5
M = 0.8                      # side margin
CW = W - 2 * M               # content width 11.733

BLACK = "AvertaStd-Black"    # SamsungSharpSans stand-in
BOLD = "AvertaStd-Bold"
BOOK = "AvertaStd-Regular"   # SamsungOne stand-in

BLUE = "1428A0"              # Samsung brand blue
AMBER = "E0A100"             # OccluSynth hero — imagined / hidden geometry
AMBER_INK = "8A6200"         # amber as text on white (contrast-safe)
FREE = "2E7D32"
SURFACE = "C0272D"
UNOBS = "6E7681"

LIGHT = dict(bg="FFFFFF", ink="000000", ink2="55595F", ink3="8B9096",
             card="F5F6F7", card2="EDEEF0", amber_ink=AMBER_INK, hair="E4E5E7")
DARK = dict(bg="0B0E14", ink="FFFFFF", ink2="A8AEB5", ink3="767C84",
            card="161B23", card2="1E242E", amber_ink=AMBER, hair="252B34")

# vertical rhythm
Y_KICKER, H_KICKER = 0.62, 0.34
Y_HEAD = 1.10
Y_DECK = 1.94
Y_BODY = 2.52                # when a deck line is present
Y_BODY_NODECK = 2.10
Y_FOOT = 6.83


def rgb(h):
    return RGBColor.from_string(h)


# ---------------------------------------------------------------- primitives
def _no_outline(shape):
    shape.line.fill.background()
    shape.shadow.inherit = False


def rect(slide, x, y, w, h, fill, radius=0.06):
    """Rounded rectangle. radius is in inches (converted to the shape adjustment)."""
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sh.adjustments[0] = min(0.5, radius / min(w, h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = rgb(fill)
    _no_outline(sh)
    return sh


def ring(slide, x, y, d, color, weight=1.25):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    sh.fill.background()
    sh.line.color.rgb = rgb(color)
    sh.line.width = Pt(weight)
    sh.shadow.inherit = False
    return sh


def dot(slide, x, y, d, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    sh.fill.solid()
    sh.fill.fore_color.rgb = rgb(color)
    _no_outline(sh)
    return sh


def dash(slide, x, y, w, color, h=0.035):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sh.adjustments[0] = 0.5
    sh.fill.solid()
    sh.fill.fore_color.rgb = rgb(color)
    _no_outline(sh)
    return sh


def tri(slide, x, y, size, color, rotation=90):
    """Small triangle — our arrow/marker glyph (Averta has no arrow or shape glyphs).
    rotation: 0 = up, 90 = right, 180 = down."""
    sh = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                Inches(x), Inches(y), Inches(size * 0.85), Inches(size))
    sh.rotation = rotation
    sh.fill.solid()
    sh.fill.fore_color.rgb = rgb(color)
    _no_outline(sh)
    return sh


def tbox(slide, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP):
    """paras: list of dicts {runs, size, font, color, bold, align, space_after,
    line_spacing, spc}. `runs` may be a plain string or a list of (text, overrides)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, spec in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = spec.get("align", PP_ALIGN.LEFT)
        if spec.get("space_after") is not None:
            p.space_after = Pt(spec["space_after"])
        if spec.get("space_before") is not None:
            p.space_before = Pt(spec["space_before"])
        if spec.get("line_spacing"):
            p.line_spacing = spec["line_spacing"]
        runs = spec["runs"]
        if isinstance(runs, str):
            runs = [(runs, {})]
        expanded = []
        for text, ov in runs:
            expanded.extend(_supsplit(text, ov))
        for text, ov in expanded:
            r = p.add_run()
            r.text = text
            f = r.font
            f.name = ov.get("font", spec.get("font", BOOK))
            f.size = Pt(ov.get("size", spec.get("size", 13)))
            f.bold = False
            f.color.rgb = rgb(ov.get("color", spec.get("color", "000000")))
            rPr = r._r.get_or_add_rPr()
            spc = ov.get("spc", spec.get("spc"))
            if spc:
                rPr.set("spc", str(int(spc * 100)))
            if ov.get("sup"):
                rPr.set("baseline", "30000")
    return tb


def _supsplit(text, ov):
    """`64{3}` -> a normal run plus a superscript run. Averta has no ³ glyph."""
    import re
    parts = re.split(r"\{([^}]*)\}", text)
    out = []
    for i, part in enumerate(parts):
        if not part:
            continue
        o = dict(ov)
        if i % 2 == 1:
            o["sup"] = True
        out.append((part, o))
    return out or [("", ov)]


def pill(slide, x, y, text, fill, fg, size=10.5, font=BOLD, pad=0.30, h=0.34, spc=0.9):
    """Samsung-style pill chip. Width is measured from the text."""
    w = measure(text, font, size) / 72.0 + 2 * pad + (spc * len(text) / 72.0)
    sh = rect(slide, x, y, w, h, fill, radius=h / 2)
    tf = sh.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.color.rgb = rgb(fg)
    r._r.get_or_add_rPr().set("spc", str(int(spc * 100)))
    return sh, w


# ---------------------------------------------------------------- measurement
_FONT_FILES = {
    BLACK: os.path.expanduser("~/Library/Fonts/AvertaStd-Black.ttf"),
    BOLD: os.path.expanduser("~/Library/Fonts/AvertaStd-Bold.ttf"),
    BOOK: os.path.expanduser("~/Library/Fonts/AvertaStd-Regular.ttf"),
}
_cache = {}
WARNINGS = []


def measure(text, font, size):
    """Width of `text` in points."""
    from PIL import ImageFont
    key = (font, round(size * 4))
    if key not in _cache:
        _cache[key] = ImageFont.truetype(_FONT_FILES[font], int(size * 4))
    return _cache[key].getlength(text) / 4.0


def wrap_lines(text, font, size, width_in):
    """Number of lines `text` needs inside width_in inches."""
    limit = width_in * 72.0
    lines, cur = 1, ""
    for word in text.split():
        trial = (cur + " " + word).strip()
        if measure(trial, font, size) <= limit or not cur:
            cur = trial
        else:
            lines += 1
            cur = word
    return lines


def fits(tag, text, font, size, width_in, height_in, leading=1.32):
    text = text.replace("{", "").replace("}", "")
    n = wrap_lines(text, font, size, width_in)
    need = n * size * leading / 72.0
    if need > height_in + 0.012:
        WARNINGS.append(f"{tag}: needs {need:.2f}in ({n} lines) in {height_in:.2f}in — "
                        f"{size}pt/{width_in:.2f}in")
    return need


# ---------------------------------------------------------------- page furniture
prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]


def new_slide(dark=False):
    s = prs.slides.add_slide(BLANK)
    t = DARK if dark else LIGHT
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = rgb(t["bg"])
    return s, t


def header(slide, t, kicker, headline, deck=None, head_size=38):
    pill(slide, M, Y_KICKER, kicker.upper(), BLUE, "FFFFFF")
    fits(f"head:{headline[:24]}", headline, BLACK, head_size, CW, 0.78)
    tbox(slide, M, Y_HEAD, CW, 0.8,
         [dict(runs=headline, font=BLACK, size=head_size, color=t["ink"])])
    if deck:
        fits(f"deck:{headline[:20]}", deck, BOOK, 14.5, CW - 1.2, 0.62)
        tbox(slide, M, Y_DECK, CW - 1.2, 0.66,
             [dict(runs=deck, font=BOOK, size=14.5, color=t["ink2"], line_spacing=1.3)])


def footnote(slide, t, text, y=Y_FOOT):
    tri(slide, M + 0.005, y + 0.055, 0.105, BLUE)
    tbox(slide, M + 0.24, y, CW - 0.24, 0.34,
         [dict(runs=text, font=BOOK, size=10.5, color=t["ink3"], line_spacing=1.25)])


def label(slide, t, x, y, w, text, color=None, size=10):
    tbox(slide, x, y, w, 0.22,
         [dict(runs=text.upper(), font=BOLD, size=size,
               color=color or t["ink3"], spc=1.2)])


# ================================================================= SLIDE 1 — title
s, t = new_slide(dark=True)
pill(s, M, 0.95, "PROBLEM STATEMENT 09  ·  FINAL PRESENTATION ROUND", BLUE, "FFFFFF")
tbox(s, M, 1.72, CW, 1.66,
     [dict(runs="OCCLUSYNTH", font=BLACK, size=92, color="FFFFFF", spc=2.2)])
tbox(s, M, 3.34, 9.6, 1.0,
     [dict(runs=[("Occlusion-aware 3D scene reconstruction in partially observable ", {}),
                 ("real-world environments", {"color": AMBER})],
           font=BOOK, size=20, color=DARK["ink2"], line_spacing=1.35)])

meta = [("Team", "onemore_adi"), ("Member", "Aditya Agarwal"),
        ("Institute", "National Institute of Technology, Rourkela"),
        ("Repository", "github.com/onemore-adi/OccluSynth")]
cols = [M, M + 2.05, M + 4.10, M + 7.70]
widths = [1.9, 1.9, 3.4, 3.9]
for (k, v), cx, cwid in zip(meta, cols, widths):
    label(s, t, cx, 5.42, cwid, k)
    fits(f"meta:{k}", v, BOOK, 13.5, cwid, 0.5)
    tbox(s, cx, 5.72, cwid, 0.52,
         [dict(runs=v, font=BOOK, size=13.5, color="FFFFFF", line_spacing=1.25)])

# ================================================================= SLIDE 2 — stakes
s, t = new_slide()
header(s, t, "The stakes", "Robots can't see behind the couch.",
       "An indoor robot sees only line-of-sight — and classical fusion forces one "
       "dishonest choice about everything else.")

# left: the 61% stat, its caption, and the market line
tbox(s, M, Y_BODY + 0.10, 3.3, 1.42,
     [dict(runs="61%", font=BLACK, size=76, color=t["ink"])])
tbox(s, M, Y_BODY + 1.46, 3.15, 1.0,
     [dict(runs="of the observable volume in a typical cluttered room is occluded — "
                "hidden behind furniture.", font=BOOK, size=13, color=t["ink2"],
           line_spacing=1.35)])
tbox(s, M, Y_BODY + 2.62, 3.15, 1.3,
     [dict(runs="Safety-critical indoor autonomy — warehouse AMRs, home service robots, "
                "inspection bots — cannot ship on “assume it’s empty.”",
           font=BOOK, size=12.5, color=t["ink2"], line_spacing=1.35)])

RX, RW = M + 3.72, CW - 3.72          # right column
fail = [("Call the unseen space FREE", SURFACE, "Silent collisions",
         "The planner drives straight through the hidden chair leg it never measured. Unsafe."),
        ("Call the unseen space BLOCKED", UNOBS, "Paralysed robot",
         "Every occluded voxel becomes a wall. The robot freezes and never reaches the goal. Useless.")]
cy = Y_BODY
for chip, ccol, title, body in fail:
    rect(s, RX, cy, RW, 1.32, t["card"], radius=0.14)
    pill(s, RX + 0.34, cy + 0.26, chip.upper(), ccol, "FFFFFF", size=9, pad=0.24, h=0.29)
    tbox(s, RX + 0.34, cy + 0.66, 3.0, 0.34,
         [dict(runs=title, font=BOLD, size=17, color=t["ink"])])
    fits(f"fail:{title}", body, BOOK, 12.5, RW - 3.78, 0.62)
    tbox(s, RX + 3.42, cy + 0.62, RW - 3.78, 0.66,
         [dict(runs=body, font=BOOK, size=12.5, color=t["ink2"], line_spacing=1.32)])
    cy += 1.50

rect(s, RX, cy + 0.06, RW, 1.06, "0B0E14", radius=0.14)
tbox(s, RX + 0.34, cy + 0.28, RW - 0.68, 0.7,
     [dict(runs=[("OccluSynth reconstructs the honest third state — ", {}),
                 ("occluded but probably occupied", {"color": AMBER, "font": BOLD}),
                 (" — so the robot routes around what it cannot see, instead of guessing.", {})],
           font=BOOK, size=13.5, color="FFFFFF", line_spacing=1.35)])
footnote(s, t, "The demo video shows this detour happening live on a real ScanNet scene — "
               "this slide is the argument behind it.")

# ================================================================= SLIDE 3 — 4 states
s, t = new_slide()
header(s, t, "Problem framing", "An honest four-state world.",
       "Classical TSDF fusion is binary: seen-solid or seen-empty. OccluSynth labels every "
       "voxel with one of four states — the distinction the planner has been missing.")

states = [(FREE, "FREE", "observed empty"),
          (SURFACE, "SURFACE", "measured solid"),
          (AMBER, "OCCLUDED", "unobserved, inferred"),
          (UNOBS, "UNOBSERVABLE", "no evidence — left alone")]
cw_, gap = 2.72, 0.29
for i, (col, name, desc) in enumerate(states):
    x = M + i * (cw_ + gap)
    hero = (name == "OCCLUDED")
    rect(s, x, Y_BODY, cw_, 1.92, t["card2"] if hero else t["card"], radius=0.16)
    rect(s, x + 0.34, Y_BODY + 0.34, 0.62, 0.62, col, radius=0.14)
    tbox(s, x + 0.34, Y_BODY + 1.10, cw_ - 0.68, 0.3,
         [dict(runs=name, font=BOLD, size=15, color=t["ink"])])
    tbox(s, x + 0.34, Y_BODY + 1.40, cw_ - 0.68, 0.4,
         [dict(runs=desc, font=BOOK, size=11.5, color=t["ink2"], line_spacing=1.25)])

rect(s, M, 4.72, CW, 1.34, t["card"], radius=0.16)
tbox(s, M + 0.42, 4.98, CW - 0.84, 0.9,
     [dict(runs=[("The key move: ", {"font": BOLD, "color": t["ink"]}),
                 ("separating ", {}),
                 ("OCCLUDED", {"font": BOLD, "color": AMBER_INK}),
                 (" (recoverable) from ", {}),
                 ("UNOBSERVABLE", {"font": BOLD, "color": t["ink"]}),
                 (" (no evidence) is what stops the system inventing geometry it cannot "
                  "justify. Every voxel reports whether it was seen or imagined.", {})],
           font=BOOK, size=14, color=t["ink2"], line_spacing=1.4)])
footnote(s, t, "The video renders these four colours filling a live room; here we define "
               "what each one licenses the model to do.")

# ================================================================= SLIDE 4 — data flow
s, t = new_slide()
header(s, t, "Data flow", "From RGB to a safe trajectory.")

io = [("INPUT", "RGB frames + sparse metric anchors", M, 5.7),
      ("OUTPUT", "Dense occlusion-aware SDF + collision-safe trajectory", M + 6.03, 5.7)]
for name, val, x, wid in io:
    rect(s, x, Y_BODY_NODECK, wid, 0.66, t["card"], radius=0.14)
    label(s, t, x + 0.3, Y_BODY_NODECK + 0.15, 1.1, name, color=t["ink3"], size=9.5)
    tbox(s, x + 1.28, Y_BODY_NODECK + 0.13, wid - 1.58, 0.42,
         [dict(runs=val, font=BOOK, size=12.5, color=t["ink"])], anchor=MSO_ANCHOR.MIDDLE)

stages = [("1", "Perception", "RGB → dense depth; RANSAC fixes metric scale"),
          ("2", "Fusion", "Visibility-aware TSDF → 4-state voxel grid"),
          ("3", "Completion", "3D U-Net writes SDF into the occluded volume"),
          ("4", "Uncertainty", "MC-dropout → per-voxel occupancy p_occ"),
          ("5", "Planner", "Risk-graded A* on the completed cost map")]
sw, sgap = 2.05, 0.37
SY = 3.30
for i, (num, name, desc) in enumerate(stages):
    x = M + i * (sw + sgap)
    hero = name in ("Completion",)
    rect(s, x, SY, sw, 2.02, t["card2"] if hero else t["card"], radius=0.16)
    dot(s, x + 0.3, SY + 0.3, 0.42, AMBER if hero else BLUE)
    tbox(s, x + 0.3, SY + 0.375, 0.42, 0.3,
         [dict(runs=num, font=BLACK, size=13,
               color="000000" if hero else "FFFFFF", align=PP_ALIGN.CENTER)])
    tbox(s, x + 0.3, SY + 0.86, sw - 0.6, 0.3,
         [dict(runs=name, font=BOLD, size=14.5, color=t["ink"])])
    desc = desc.replace("→", "to")
    fits(f"stage:{name}", desc, BOOK, 11.5, sw - 0.6, 0.85)
    tbox(s, x + 0.3, SY + 1.18, sw - 0.6, 0.85,
         [dict(runs=desc, font=BOOK, size=11.5, color=t["ink2"], line_spacing=1.28)])
    if i < len(stages) - 1:
        tri(s, x + sw + 0.135, SY + 0.93, 0.13, t["ink3"])

rect(s, M, 5.62, CW, 1.02, "0B0E14", radius=0.14)
tbox(s, M + 0.42, 5.83, CW - 0.84, 0.64,
     [dict(runs=[("The visibility mask is conditioning", {"font": BOLD, "color": AMBER}),
                 (" — it tells the completer exactly which voxels it is licensed to imagine, "
                  "and which to leave untouched. Generation is fenced to the recoverable region.",
                  {})],
           font=BOOK, size=13.5, color="FFFFFF", line_spacing=1.35)])
footnote(s, t, "The video animates voxels filling in; this diagram is the input/output "
               "contract a moving image can’t freeze.", y=6.85)

# ================================================================= SLIDE 5 — architecture
s, t = new_slide()
header(s, t, "Architecture", "The decision behind each stage.")

rows = [("A", "VGGT-Omega, frozen",
         "Preserve the CVPR’25 geometric prior; RANSAC fixes metric scale (ARE 2.4%) — "
         "cheaper and safer than fine-tuning a foundation model."),
        ("B", "4-state TSDF, not binary",
         "Binary occupancy hides the recoverable-vs-no-evidence gap. The extra state is "
         "exactly what the planner needs to act honestly."),
        ("C", "Masked L1 on occluded only",
         "Supervise imagination precisely where it is licensed; never penalise observed "
         "surface. 14.7 M params, GT SDF from meshes."),
        ("D", "MC-dropout uncertainty",
         "A point estimate is not enough — the planner needs per-voxel confidence (p_occ) "
         "to grade risk, not just a guess."),
        ("E", "Risk-graded A*, soft cost",
         "Occluded space is graded, not an ∞ wall. Honesty as a soft cost lets the robot "
         "keep moving while respecting hidden risk.")]
label(s, t, M + 0.88, Y_BODY_NODECK - 0.16, 3.0, "Decision")
label(s, t, M + 4.20, Y_BODY_NODECK - 0.16, 3.0, "Why")
ry, rh, rgap = Y_BODY_NODECK + 0.12, 0.80, 0.13
for letter, title, why in rows:
    rect(s, M, ry, CW, rh, t["card"], radius=0.14)
    dot(s, M + 0.30, ry + 0.20, 0.40, BLUE)
    tbox(s, M + 0.30, ry + 0.275, 0.40, 0.28,
         [dict(runs=letter, font=BLACK, size=12.5, color="FFFFFF", align=PP_ALIGN.CENTER)])
    tbox(s, M + 0.88, ry + 0.27, 3.05, 0.32,
         [dict(runs=title, font=BOLD, size=15, color=t["ink"])])
    fits(f"why:{letter}", why, BOOK, 12, CW - 4.55, 0.52)
    tbox(s, M + 4.20, ry + 0.16, CW - 4.55, 0.56,
         [dict(runs=why, font=BOOK, size=12, color=t["ink2"], line_spacing=1.32)],
         anchor=MSO_ANCHOR.MIDDLE)
    ry += rh + rgap
footnote(s, t, "The video shows the pipeline running end-to-end; this slide is the "
               "reasoning the demo can’t narrate.", y=6.86)

# ================================================================= SLIDE 6 — results (dark)
s, t = new_slide(dark=True)
header(s, t, "From problem to result", "How it solves the problem.",
       "The third state turns “unknown” into a graded cost the planner can act on — and we "
       "measure whether it was right on a purpose-built safety benchmark.")

stats = [("Hidden-hazard awareness", "baseline  0%", "21%",
          "line-of-sight detects none of them"),
         ("Occluded F-score @ 5 cm", "baseline  0%", "32%",
          "TSDF cannot see behind walls"),
         ("Occluded MAE", "from  45.3 cm", "27.1 cm",
          "depth error, halved (lower is better)")]
scw, sgp = 3.71, 0.30
for i, (lab, base, big, note) in enumerate(stats):
    x = M + i * (scw + sgp)
    rect(s, x, Y_BODY, scw, 2.16, t["card"], radius=0.16)
    label(s, t, x + 0.36, Y_BODY + 0.32, scw - 0.72, lab, color=t["ink3"], size=9.5)
    tbox(s, x + 0.36, Y_BODY + 0.60, scw - 0.72, 0.28,
         [dict(runs=base, font=BOOK, size=12, color=t["ink3"])])
    fits(f"stat:{lab}", big, BLACK, 44, scw - 0.72, 0.80)
    tbox(s, x + 0.36, Y_BODY + 0.88, scw - 0.72, 0.80,
         [dict(runs=big, font=BLACK, size=44, color=AMBER)])
    tbox(s, x + 0.36, Y_BODY + 1.74, scw - 0.72, 0.34,
         [dict(runs=note, font=BOOK, size=11.5, color=t["ink2"], line_spacing=1.25)])

tbox(s, M, 5.10, CW, 1.1,
     [dict(runs=[("Silent collisions and the frozen robot both disappear when occluded space "
                  "is completed and priced by confidence. On scene0556_00 the planner achieves ",
                  {}),
                 ("15.5% collision-avoidance", {"font": BOLD, "color": "FFFFFF"}),
                 (" over a line-of-sight baseline, detouring around inferred hazard clusters "
                  "on a 13.6 m path.", {})],
           font=BOOK, size=14, color=t["ink2"], line_spacing=1.45)])
footnote(s, t, "The demo shows one detour; these numbers show it holds across the benchmark, "
               "not just the clip.")

# ================================================================= SLIDE 7 — novelty
s, t = new_slide()
header(s, t, "Novelty", "Knows what it doesn’t know.",
       "The first system that reconstructs occluded geometry, reports its own confidence, "
       "and measures whether it was right. Capability coverage vs. the closest published systems:")

col_x = [M, 5.30, 6.85, 8.40, 9.95, 11.50]
col_w = [4.30, 1.45, 1.45, 1.45, 1.45, 1.53]
heads = ["Capability", "Atlas /\nNeuralRecon", "VGGT", "3D Diffusion\n(DiffInDScene)",
         "Behind\nthe Veil", "OccluSynth\n(ours)"]
HY = Y_BODY + 0.02
rect(s, col_x[5] - 0.12, HY - 0.14, col_w[5] + 0.22, 3.46, t["card2"], radius=0.14)
for i, htxt in enumerate(heads):
    align = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER
    is_ours = (i == 5)
    tbox(s, col_x[i], HY, col_w[i], 0.62,
         [dict(runs=line, font=BOLD if is_ours else BOOK, size=11,
               color=t["ink"] if is_ours else t["ink3"], align=align, line_spacing=1.2)
          for line in htxt.split("\n")], anchor=MSO_ANCHOR.BOTTOM)

caps = [("Reconstructs occluded geometry", ["no", "no", "part", "part", "yes"]),
        ("Works from sparse metric anchors", ["no", "no", "no", "no", "yes"]),
        ("Uncertainty-aware completion", ["no", "part", "no", "no", "yes"]),
        ("Measured on an occlusion safety benchmark", ["no", "no", "no", "no", "yes"]),
        ("Closed-loop into a navigation planner", ["no", "no", "no", "no", "yes"])]
ry = HY + 0.78
for cap, marks in caps:
    tbox(s, col_x[0], ry - 0.02, col_w[0], 0.34,
         [dict(runs=cap, font=BOOK, size=12.5, color=t["ink"])])
    for j, mk in enumerate(marks):
        cx = col_x[j + 1] + col_w[j + 1] / 2
        ours = (j == 4)
        if mk == "yes":
            dot(s, cx - 0.085, ry, 0.17, BLUE if ours else t["ink"])
        elif mk == "part":
            ring(s, cx - 0.085, ry, 0.17, t["ink3"])
        else:
            dash(s, cx - 0.09, ry + 0.068, 0.18, t["hair"])
    ry += 0.54

LGY = ry + 0.28
lg = [("yes", "yes"), ("part", "partial"), ("no", "no")]
lx = M
for mk, lbl in lg:
    if mk == "yes":
        dot(s, lx, LGY, 0.15, BLUE)
    elif mk == "part":
        ring(s, lx, LGY, 0.15, t["ink3"])
    else:
        dash(s, lx, LGY + 0.055, 0.16, t["hair"])
    tbox(s, lx + 0.26, LGY - 0.025, 0.8, 0.26,
         [dict(runs=lbl, font=BOOK, size=11, color=t["ink3"])])
    lx += 1.05
tbox(s, M, LGY + 0.52, CW, 0.5,
     [dict(runs="Semantic labelling (5-class) is deferred to Phase 2 and intentionally "
                "excluded from the MVP claims above.",
           font=BOOK, size=11.5, color=t["ink3"])])

# ================================================================= SLIDE 8 — results tables
s, t = new_slide()
header(s, t, "Evidence", "Results & validation.")

mini = [("Occluded MAE", "45.3", "27.1 cm", "lower is better"),
        ("Occluded F@5cm", "0%", "32%", "TSDF cannot see behind walls"),
        ("Hidden-hazard awareness", "0%", "21%", "line-of-sight detects none")]
mw, mg = 3.71, 0.30
for i, (lab, before, after, note) in enumerate(mini):
    x = M + i * (mw + mg)
    rect(s, x, Y_BODY_NODECK, mw, 1.16, t["card"], radius=0.14)
    label(s, t, x + 0.32, Y_BODY_NODECK + 0.21, mw - 0.64, lab, size=9.5)
    bx = x + 0.32
    bw = measure(before, BLACK, 22) / 72.0
    tbox(s, bx, Y_BODY_NODECK + 0.48, bw + 0.1, 0.42,
         [dict(runs=before, font=BLACK, size=22, color=t["ink3"])])
    tri(s, bx + bw + 0.15, Y_BODY_NODECK + 0.575, 0.13, t["ink3"])
    tbox(s, bx + bw + 0.44, Y_BODY_NODECK + 0.48, mw - 0.76 - bw, 0.42,
         [dict(runs=after, font=BLACK, size=22, color=AMBER_INK)])
    tbox(s, x + 0.32, Y_BODY_NODECK + 0.86, mw - 0.64, 0.26,
         [dict(runs=note, font=BOOK, size=10.5, color=t["ink2"])])

TY = 3.62
tables = [
    ("Completer — occluded voxels",
     [("Method", None), ("MAE (cm)", "down"), ("Sign acc", "up"), ("Compl<5cm", "up")],
     [["no_completion", "45.27", "0.299", "0.061"],
      ["occluded_as_free", "42.00", "0.701", "0.121"],
      ["OccluSynth Completer", "27.14", "0.722", "0.349"]], M, 5.70),
    ("Geometry — surface vs. occluded",
     [("Method", None), ("Chamfer-L1", "down"), ("F@5cm", "up"), ("Occl. F", "up")],
     [["TSDF-only", "3.11 cm", "74.1%", "0.0%"],
      ["OccluSynth", "1.77 cm", "83.5%", "32.0%"]], M + 6.03, 5.70)]
for title, cols_, rows_, x, wid in tables:
    label(s, t, x, TY, wid, title, color=t["ink3"], size=9.5)
    colw = [wid * 0.40, wid * 0.20, wid * 0.20, wid * 0.20]
    cx = [x + sum(colw[:i]) for i in range(4)]
    for i, (c, direction) in enumerate(cols_):
        pad = 0.16 if direction else 0.0
        tbox(s, cx[i], TY + 0.34, colw[i] - pad, 0.3,
             [dict(runs=c, font=BOOK, size=10.5, color=t["ink3"],
                   align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.RIGHT)])
        if direction:
            tri(s, cx[i] + colw[i] - 0.115, TY + 0.395, 0.10, t["ink3"],
                rotation=0 if direction == "up" else 180)
    ry = TY + 0.70
    for r in rows_:
        best = r[0].startswith("OccluSynth")
        if best:
            rect(s, x - 0.16, ry - 0.06, wid + 0.32, 0.46, t["card2"], radius=0.10)
        for i, cell in enumerate(r):
            tbox(s, cx[i], ry + 0.04, colw[i], 0.32,
                 [dict(runs=cell, font=BOLD if best else BOOK, size=12.5,
                       color=t["ink"] if best else t["ink2"],
                       align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.RIGHT)])
        ry += 0.50

tbox(s, M, 5.95, CW, 0.78,
     [dict(runs=[("Also validated:  ", {"font": BOLD, "color": t["ink"]}),
                 ("metric grounding 2.4% ARE across 10/10 val scenes · 61% of observable "
                  "volume occluded on scene0000_00 · 55+ automated tests, 18 planner tests. "
                  "All metrics from the interim 64{3} MPS checkpoint (epoch 32) — the scripted "
                  "96{3} A100 run is expected to lift every number.", {})],
           font=BOOK, size=11.5, color=t["ink2"], line_spacing=1.35)])

# ================================================================= SLIDE 9 — OSS
s, t = new_slide()
header(s, t, "OSS + models", "Open-source foundations.")

groups = [("Reused — frozen", UNOBS,
           [("VGGT-Omega",
             "Meta AI feed-forward 3D geometry transformer (CVPR’25 Best Paper). RGB → "
             "depth/pose, used as-is; scale fixed downstream by RANSAC."),
            ("ScanNet v2 · 7-Scenes",
             "1,513 indoor RGB-D scenes (supervision + metrics); 7-Scenes as an unchanged "
             "cross-dataset portability check.")]),
          ("Trained & publishing", BLUE,
           [("OccluSynth Completer",
             "3D U-Net, 14.7 M params, masked L1 on surface + occluded voxels. Predicts SDF "
             "inside the occluded volume."),
            ("Checkpoint on Hugging Face",
             "interim_64_aug, epoch 32. Open weights + reproducible eval harness; HF publish "
             "in progress.")]),
          ("Phase 2", AMBER,
           [("3D Latent Diffusion",
             "DiffInDScene-family indoor prior to replace the U-Net completer with "
             "visibility-aware conditioning."),
            ("Scoped post-hackathon",
             "A stronger generative prior for the same fenced occluded region — drop-in "
             "behind the existing mask.")])]
gw, ggap = 3.71, 0.30
for i, (gname, gcol, items) in enumerate(groups):
    x = M + i * (gw + ggap)
    pill(s, x, Y_BODY_NODECK, gname.upper(), gcol,
         "000000" if gcol == AMBER else "FFFFFF", size=9, pad=0.24, h=0.30)
    iy = Y_BODY_NODECK + 0.50
    for iname, ibody in items:
        rect(s, x, iy, gw, 1.72, t["card"], radius=0.16)
        fits(f"oss:{iname}", iname, BOLD, 14.5, gw - 0.64, 0.62)
        tbox(s, x + 0.32, iy + 0.28, gw - 0.64, 0.62,
             [dict(runs=iname, font=BOLD, size=14.5, color=t["ink"], line_spacing=1.2)])
        body = ibody.replace("→", "to")
        fits(f"ossb:{iname}", body, BOOK, 11.5, gw - 0.64, 0.86)
        tbox(s, x + 0.32, iy + 0.86, gw - 0.64, 0.86,
             [dict(runs=body, font=BOOK, size=11.5, color=t["ink2"], line_spacing=1.3)])
        iy += 1.90

rect(s, M, 6.28, CW, 0.80, "0B0E14", radius=0.14)
tbox(s, M + 0.40, 6.46, CW - 0.80, 0.46,
     [dict(runs=[("New work in this project:  ", {"font": BOLD, "color": AMBER}),
                 ("visibility-aware TSDF · RANSAC metric grounding · the completer · the "
                  "risk-graded A* planner · the full evaluation harness · the first "
                  "reproducible occlusion-safety benchmark — no simulator required.", {})],
           font=BOOK, size=12, color="FFFFFF", line_spacing=1.3)])

# ================================================================= SLIDE 10 — deployment
s, t = new_slide()
header(s, t, "Real-world usage", "Routing around the unseen.")

LW = 5.70
paths = [("Line-of-sight planner", SURFACE, "unseen hazard",
          "Path clips the obstacle it never measured — collision.", False),
         ("OccluSynth planner", FREE, "inferred occupied",
          "Detours around the inferred region — 15.5% fewer collisions.", True)]
py = Y_BODY_NODECK
for pname, pcol, tag, desc, good in paths:
    rect(s, M, py, LW, 1.68, t["card"], radius=0.16)
    tbox(s, M + 0.34, py + 0.26, LW - 0.68, 0.32,
         [dict(runs=pname, font=BOLD, size=15, color=t["ink"])])
    pill(s, M + 0.34, py + 0.66, tag.upper(), AMBER if good else t["hair"],
         "000000" if good else t["ink2"], size=8.5, pad=0.22, h=0.27)
    dot(s, M + 0.34, py + 1.14, 0.17, pcol)
    fits(f"path:{pname}", desc, BOOK, 12, LW - 1.05, 0.4)
    tbox(s, M + 0.62, py + 1.10, LW - 1.0, 0.4,
         [dict(runs=desc, font=BOOK, size=12, color=t["ink2"], line_spacing=1.3)])
    py += 1.86

RX2 = M + 6.03
label(s, t, RX2, Y_BODY_NODECK + 0.04, 5.7, "Who needs this")
who = [("Warehouse AMRs", "moving fast through aisles blocked by pallets and shelving."),
       ("Home service robots", "navigating rooms where furniture hides most of the floor."),
       ("Inspection & delivery bots", "where “assume it’s empty” is a safety incident, not a bug.")]
wy = Y_BODY_NODECK + 0.34
for wname, wdesc in who:
    dot(s, RX2, wy + 0.09, 0.13, BLUE)
    tbox(s, RX2 + 0.26, wy, 5.44, 0.50,
         [dict(runs=[(wname + " — ", {"font": BOLD, "color": t["ink"]}), (wdesc, {})],
               font=BOOK, size=12, color=t["ink2"], line_spacing=1.3)])
    wy += 0.52

label(s, t, RX2, wy + 0.18, 5.7, "Path to production")
prod = [("Poses", "Swap ScanNet GT poses for a VIO / SLAM front-end (real-time on-robot)."),
        ("Weights", "Publish the completer checkpoint to Hugging Face — open, reproducible."),
        ("Fidelity", "Run the scripted 96{3} A100 training for production-grade completion."),
        ("Calibration", "Temperature-scale uncertainty to ECE < 0.05 before closing the loop.")]
ppy = wy + 0.50
for i, (pn, pd) in enumerate(prod, 1):
    tbox(s, RX2, ppy, 0.3, 0.3,
         [dict(runs=str(i), font=BLACK, size=11.5, color=t["ink3"])])
    tbox(s, RX2 + 0.30, ppy, 5.40, 0.46,
         [dict(runs=[(pn + " — ", {"font": BOLD, "color": t["ink"]}), (pd, {})],
               font=BOOK, size=11.5, color=t["ink2"], line_spacing=1.3)])
    ppy += 0.44
footnote(s, t, "The video is the live proof-of-concept; this slide is the market and the "
               "road from demo to robot.", y=6.9)

# ================================================================= SLIDE 11 — status (dark)
s, t = new_slide(dark=True)
header(s, t, "Close", "Status, honest gaps & roadmap.")

cols3 = [("Built & tested", FREE,
          ["Metric grounding — RANSAC, ARE 2.4%",
           "Visibility-aware 4-state TSDF",
           "3D U-Net completer — beats both baselines",
           "Risk-graded A* planner",
           "MC-dropout uncertainty",
           "ScanNet-native occlusion safety benchmark",
           "55+ automated tests across the pipeline"]),
         ("Honest gaps", AMBER,
          ["Uncertainty uncalibrated — ECE 0.42 (fix scoped)",
           "Completer at interim 64{3} on Apple MPS",
           "Cross-dataset probe uses 7-Scenes metric depth",
           "Checkpoint not yet published to Hugging Face",
           "Poses are ScanNet GT (VIO in deployment)"]),
         ("Roadmap — Phase 2", BLUE,
          ["Full 96{3} A100 run — lifts every metric",
           "Temperature scaling → ECE < 0.05",
           "3D diffusion completer replaces U-Net",
           "5-class semantic head (mIoU)",
           "Active perception — next-best-view"])]
c3w, c3g = 3.71, 0.30
for i, (cname, ccol, items) in enumerate(cols3):
    x = M + i * (c3w + c3g)
    rect(s, x, Y_BODY_NODECK, c3w, 3.94, t["card"], radius=0.16)
    pill(s, x + 0.30, Y_BODY_NODECK + 0.26, cname.upper(), ccol,
         "000000" if ccol == AMBER else "FFFFFF", size=8.5, pad=0.22, h=0.28)
    iy = Y_BODY_NODECK + 0.74
    ITW = c3w - 0.84
    for it in items:
        dot(s, x + 0.32, iy + 0.075, 0.10, ccol)
        it = it.replace("→", "to")
        # count lines against a 6%-narrower box so a slightly different PowerPoint
        # line-break can't collide the next item into this one
        n = wrap_lines(it.replace("{", "").replace("}", ""), BOOK, 10.5, ITW * 0.97)
        tbox(s, x + 0.52, iy, ITW, 0.19 * n + 0.08,
             [dict(runs=it, font=BOOK, size=10.5, color=t["ink2"], line_spacing=1.25)])
        iy += 0.17 + 0.19 * n

tbox(s, M, 6.32, CW, 0.8,
     [dict(runs=[("Compute:  ", {"font": BOLD, "color": "FFFFFF"}),
                 ("interim training on Apple Silicon MPS (16 GB) — no cloud GPU used. The "
                  "full 96{3} run is scripted & data-prepared (train_completer.py --device cuda "
                  "--crop_size 96); ~100 A100 GPU-hr away. ", {}),
                 ("This is a compute gap, not missing work.", {"font": BOLD, "color": AMBER})],
           font=BOOK, size=12, color=t["ink2"], line_spacing=1.4)])

# ================================================================= SLIDE 12 — appendix
s, t = new_slide()
header(s, t, "Reference", "Appendix — references & methods.")

label(s, t, M, Y_BODY_NODECK + 0.04, 5.7, "State of the art referenced")
refs = [("Reconstruction", "Atlas (ECCV’20) · NeuralRecon (CVPR’21) · TransformerFusion "
                           "(NeurIPS’21) · VGGT (CVPR’25, our backbone)"),
        ("Completion / generation", "DiffInDScene · Octree Latent Semantic Diffusion (’25) · "
                                    "Behind the Veil (’24) · RecGen (’25)"),
        ("Depth & planning", "Marigold-DC · normal-guided sparse sampling · OA-MPC · RAMP "
                             "(risk-aware mapping & planning)")]
ry = Y_BODY_NODECK + 0.36
for rname, rbody in refs:
    fits(f"ref:{rname}", rname + " — " + rbody, BOOK, 12, 5.55, 0.78)
    tbox(s, M, ry, 5.55, 0.8,
         [dict(runs=[(rname + " — ", {"font": BOLD, "color": t["ink"]}), (rbody, {})],
               font=BOOK, size=12, color=t["ink2"], line_spacing=1.35)])
    ry += 0.86

RX3 = M + 6.03
label(s, t, RX3, Y_BODY_NODECK + 0.04, 5.7, "Implementation details")
impl = [("Metric grounding", "per-frame RANSAC affine fit on ≈500 stratified anchors; "
                             "ARE 2.4%, scale varies per scene."),
        ("Visibility TSDF", "projective integration, obliquity-corrected surface band; "
                            "per-voxel (sdf, weight, p_observed); 4 states."),
        ("Completer", "3D U-Net (14.7 M), masked L1 on surface + occluded, D4 augmentation, "
                      "96{3} crops; GT SDF from meshes."),
        ("Planner", "A* on a 2D cost map at robot height; free 0 / surface ∞ / occluded "
                    "graded by p_occupied."),
        ("Safety benchmark", "hidden hazard = occluded and GT-occupied; awareness + "
                             "collision-avoidance vs. baselines."),
        ("Metrics", "Chamfer-L1, F-score@5cm, sign accuracy, ECE — reported separately for "
                    "surface and occluded voxels.")]
iy = Y_BODY_NODECK + 0.36
for iname, ibody in impl:
    n = wrap_lines(iname + " — " + ibody, BOOK, 11.5, 5.55)
    fits(f"impl:{iname}", iname + " — " + ibody, BOOK, 11.5, 5.55, 0.62)
    tbox(s, RX3, iy, 5.55, 0.62,
         [dict(runs=[(iname + " — ", {"font": BOLD, "color": t["ink"]}), (ibody, {})],
               font=BOOK, size=11.5, color=t["ink2"], line_spacing=1.32)])
    iy += 0.20 + 0.21 * n

rect(s, M, 6.30, CW, 0.62, t["card"], radius=0.13)
tbox(s, M + 0.36, 6.47, CW - 0.72, 0.3,
     [dict(runs="Diffusion completer and 5-class semantics are Phase 2.",
           font=BOOK, size=11.5, color=t["ink2"])], anchor=MSO_ANCHOR.MIDDLE)

# ---------------------------------------------------------------- save
prs.save(OUT)
print(f"saved {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
if WARNINGS:
    print(f"\n{len(WARNINGS)} fit warnings:")
    for w in WARNINGS:
        print("  !", w)
else:
    print("no fit warnings")
